import streamlit as st
import pandas as pd
from agents import generate_dashboard
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from pptx import Presentation
import io
import matplotlib.pyplot as plt
import re
import streamlit.components.v1 as components

# ---------- CONFIG ----------
st.set_page_config(page_title="AI Dashboard", layout="wide")

# ---------- STYLE ----------
st.markdown("""
<style>
.big-title {
    font-size: 40px;
    font-weight: bold;
}
.stTextInput input {
    height: 50px;
    font-size: 18px;
}
</style>
""", unsafe_allow_html=True)

# ---------- TITLE ----------
st.markdown('<div class="big-title">🤖 Autonomous Report Generator</div>', unsafe_allow_html=True)

st.markdown("### Enter your topic")

# ---------- SESSION ----------
if "topic" not in st.session_state:
    st.session_state.topic = ""

# ---------- INPUT ----------
topic = st.text_input("", value=st.session_state.topic, key="topic_input")

# ---------- VOICE INPUT (REAL MIC) ----------
voice_html = """
<button onclick="startDictation()" style="
padding:10px 20px;
border-radius:10px;
border:1px solid #ccc;
background-color:#111827;
color:white;
font-size:16px;">
🎤 Voice Input
</button>

<script>
function startDictation() {

    var recognition = new webkitSpeechRecognition();
    recognition.lang = "en-US";

    recognition.onresult = function(event) {
        var text = event.results[0][0].transcript;

        window.parent.postMessage({
            type: "streamlit:setComponentValue",
            value: text
        }, "*");
    };

    recognition.start();
}
</script>
"""

spoken_text = components.html(voice_html, height=80)

# ---------- UPDATE INPUT FROM VOICE ----------
if spoken_text:
    topic = spoken_text
    st.session_state.topic = topic

# ---------- PDF ----------
def create_pdf(data, topic):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()

    content = []
    content.append(Paragraph(f"<b>Report: {topic}</b>", styles["Title"]))
    content.append(Spacer(1, 20))

    for i, section in enumerate(data, start=1):
        content.append(Paragraph(f"<b>{i}. {section['title']}</b>", styles["Heading2"]))
        content.append(Spacer(1, 10))

        content.append(Paragraph(section["content"], styles["Normal"]))
        content.append(Spacer(1, 10))

        for p in section["points"]:
            content.append(Paragraph(f"• {p}", styles["Normal"]))
            content.append(Spacer(1, 5))

        content.append(Spacer(1, 20))

    doc.build(content)
    buffer.seek(0)
    return buffer

# ---------- PPT ----------
def create_ppt(data):
    prs = Presentation()

    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "AI Report"
    slide.placeholders[1].text = "Generated using Agentic AI"

    for section in data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = section["title"]
        points = section["points"][:4]
        content = "\n".join([f"• {p}" for p in points])
        slide.placeholders[1].text = content

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# ---------- GENERATE ----------
if st.button("Generate Report", key="generate_btn"):

    if topic:
        st.session_state.topic = topic

        data = generate_dashboard(topic)

        st.success("Report Generated!")

        # ---------- DISPLAY ----------
        st.markdown(f"# 📊 Report: {topic}")

        for i, section in enumerate(data, start=1):
            st.markdown(f"## {i}. {section['title']}")
            st.write(section["content"])

            for idx, p in enumerate(section["points"], start=1):
                st.write(f"{idx}. {p}")

        # ---------- CHART ----------
        st.markdown("## 📈 Analysis")

        all_stats = {}
        for section in data:
            stats = section.get("stats", {})
            for k, v in stats.items():
                try:
                    all_stats[k] = int(v)
                except:
                    all_stats[k] = 50

        df = pd.DataFrame(list(all_stats.items()), columns=["Metric", "Value"])

        if not df.empty:
            st.bar_chart(df.set_index("Metric"))
            st.line_chart(df.set_index("Metric"))

            fig, ax = plt.subplots()
            df.set_index("Metric").plot.pie(y="Value", autopct="%1.1f%%", ax=ax)
            st.pyplot(fig)
        else:
            st.warning("No chart data available")

        # ---------- DOWNLOAD ----------
        clean_topic = re.sub(r'[^a-zA-Z0-9 ]', '', topic)
        pdf_name = clean_topic.replace(" ", "_") + ".pdf"
        ppt_name = clean_topic.replace(" ", "_") + ".pptx"

        col1, col2 = st.columns(2)

        with col1:
            pdf_file = create_pdf(data, topic)
            st.download_button("📄 Download PDF", pdf_file, file_name=pdf_name)

        with col2:
            ppt_file = create_ppt(data)
            st.download_button("📊 Download PPT", ppt_file, file_name=ppt_name)

    else:
        st.warning("Please speak or enter a topic")
